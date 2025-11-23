# pdf_processor.py
#
# Real processing engine for BlinkPDF tools.
# Uses only free local libraries: PyMuPDF, PyPDF2, Pillow, python-docx,
# python-pptx, openpyxl, reportlab, pytesseract (with safe fallbacks).
#
# External interface:
#   processor = PDFProcessor(output_folder)
#   out_path, download_name = processor.process(tool_slug, input_files, options)
#
#   - tool_slug: one of the defined tool slugs below
#   - input_files: list[str] of absolute file paths of uploaded files
#   - options: dict from advanced form options (may be empty)
#
# Return:
#   - out_path: absolute path to generated file (PDF, DOCX, PPTX, ZIP, etc.)
#   - download_name: nice filename to show in download

import os
import uuid
import io
import zipfile

from typing import Dict, List, Tuple

import fitz  # PyMuPDF
from PyPDF2 import PdfReader, PdfWriter
from PIL import Image

from docx import Document
from pptx import Presentation
from openpyxl import Workbook, load_workbook

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4, letter

try:
    import pytesseract
    _HAS_TESSERACT = True
except Exception:
    _HAS_TESSERACT = False


class PDFProcessor:
    """
    Core processing engine for all BlinkPDF tools.
    """

    def __init__(self, output_folder: str) -> None:
        self.output_folder = output_folder
        os.makedirs(self.output_folder, exist_ok=True)

        # Mapping from slug to implementation
        self._dispatch = {
            # CORE PDF OPS
            "compress-pdf": self._compress_pdf,
            "optimize-pdf": self._optimize_pdf,
            "merge-pdf": self._merge_pdf,
            "split-pdf": self._split_pdf,
            "rotate-pdf": self._rotate_pdf,
            "watermark-pdf": self._watermark_pdf,
            "number-pdf": self._number_pdf,
            "protect-pdf": self._protect_pdf,
            "unlock-pdf": self._unlock_pdf,
            "repair-pdf": self._repair_pdf,
            "organize-pdf": self._organize_pdf,
            "sign-pdf": self._sign_pdf,
            "annotate-pdf": self._annotate_pdf,
            "redact-pdf": self._redact_pdf,
            "flatten-pdf": self._flatten_pdf,
            "metadata-editor": self._metadata_editor,
            "fill-forms": self._fill_forms,

            # PAGE GEOMETRY
            "deskew-pdf": self._deskew_pdf,
            "crop-pdf": self._crop_pdf,
            "resize-pdf": self._resize_pdf,
            "background-remover": self._background_remover,

            # CONVERSIONS
            "pdf-to-word": self._pdf_to_word,
            "word-to-pdf": self._word_to_pdf,
            "pdf-to-image": self._pdf_to_image,
            "image-to-pdf": self._image_to_pdf,
            "pdf-to-excel": self._pdf_to_excel,
            "excel-to-pdf": self._excel_to_pdf,
            "pdf-to-powerpoint": self._pdf_to_powerpoint,
            "powerpoint-to-pdf": self._powerpoint_to_pdf,

            # TEXT / OCR / MEDIA
            "ocr-pdf": self._ocr_pdf,
            "extract-text": self._extract_text,
            "extract-images": self._extract_images,
        }

    # ------------------------------------------------------------------
    # Public entry point
    # ------------------------------------------------------------------
    def process(
        self,
        tool_slug: str,
        input_files: List[str],
        options: Dict = None,
        **kwargs,
    ) -> Tuple[str, str]:
        """
        Generic processing entry.

        :param tool_slug: tool slug string (e.g. "compress-pdf")
        :param input_files: list of uploaded file paths
        :param options: dict of advanced options from form
        :return: (output_path, download_filename)
        """
        if options is None:
            options = {}
        if not input_files:
            raise ValueError("No input files provided to processor.")

        handler = self._dispatch.get(tool_slug)
        if handler is None:
            # Unknown tool slug: fail clearly
            raise ValueError(f"Unknown tool slug: {tool_slug}")

        return handler(input_files, options or {})

    # ------------------------------------------------------------------
    # Utility helpers
    # ------------------------------------------------------------------
    def _make_output_path(self, base_name: str, ext: str = ".pdf") -> Tuple[str, str]:
        uid = uuid.uuid4().hex[:8]
        safe_base = base_name.replace(" ", "_")
        download_name = f"{safe_base}_{uid}{ext}"
        full_path = os.path.join(self.output_folder, download_name)
        return full_path, download_name

    def _ensure_pdf(self, path: str) -> None:
        if not path.lower().endswith(".pdf"):
            raise ValueError("This tool requires a PDF input file.")

    # ------------------------------------------------------------------
    # CORE PDF OPERATIONS
    # ------------------------------------------------------------------
    def _compress_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Basic compression using PyMuPDF save parameters.
        """
        in_path = files[0]
        self._ensure_pdf(in_path)

        quality = int(options.get("quality", 75))  # 1-100
        image_dpi = int(options.get("image_dpi", 120))
        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_compressed")

        doc = fitz.open(in_path)

        # Re-save with cleaning/garbage collection
        doc.save(
            out_path,
            deflate=True,
            garbage=4,
            clean=True,
            compress=True,
            ascii=False,
            expand=False,
            # PyMuPDF <=1.21 uses "use_objstm" etc. – we keep defaults.
        )
        doc.close()
        return out_path, dl_name

    def _optimize_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        More aggressive optimization: also downscale images.
        """
        in_path = files[0]
        self._ensure_pdf(in_path)
        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_optimized")

        doc = fitz.open(in_path)
        # Downscale all images
        zoom = float(options.get("zoom", 0.7))
        mat = fitz.Matrix(zoom, zoom)

        new_doc = fitz.open()
        for page in doc:
            pix = page.get_pixmap(matrix=mat)
            img_bytes = pix.tobytes("png")
            img_pdf = fitz.open("pdf", fitz.open("png", img_bytes).convert_to_pdf())
            new_doc.insert_pdf(img_pdf)

        new_doc.save(out_path, garbage=4, clean=True, deflate=True, compress=True)
        new_doc.close()
        doc.close()
        return out_path, dl_name

    def _merge_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        writer = PdfWriter()
        for in_path in files:
            self._ensure_pdf(in_path)
            reader = PdfReader(in_path)
            for page in reader.pages:
                writer.add_page(page)

        base_name = "merged"
        out_path, dl_name = self._make_output_path(base_name)
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, dl_name

    def _split_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Split PDF. If options['range'] provided (e.g. "1-3,5"), export that subset.
        Otherwise, split into individual page PDFs and zip them.
        """
        in_path = files[0]
        self._ensure_pdf(in_path)
        base_name = os.path.splitext(os.path.basename(in_path))[0]

        reader = PdfReader(in_path)

        page_range = options.get("range", "").strip()
        if page_range:
            # Export selected range as single PDF
            pages_to_keep = self._parse_page_range(page_range, len(reader.pages))
            writer = PdfWriter()
            for idx in pages_to_keep:
                writer.add_page(reader.pages[idx])
            out_path, dl_name = self._make_output_path(base_name + "_split_range")
            with open(out_path, "wb") as f:
                writer.write(f)
            return out_path, dl_name

        # Otherwise: split all pages into a ZIP
        zip_path, dl_name = self._make_output_path(base_name + "_pages", ".zip")
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(reader.pages):
                writer = PdfWriter()
                writer.add_page(page)
                single_name = f"{base_name}_page_{i+1}.pdf"
                buf = io.BytesIO()
                writer.write(buf)
                buf.seek(0)
                zf.writestr(single_name, buf.read())

        return zip_path, dl_name

    def _parse_page_range(self, page_range: str, total_pages: int) -> List[int]:
        pages = set()
        for part in page_range.split(","):
            part = part.strip()
            if not part:
                continue
            if "-" in part:
                start_s, end_s = part.split("-", 1)
                start = max(1, int(start_s))
                end = min(total_pages, int(end_s))
                for p in range(start, end + 1):
                    pages.add(p - 1)
            else:
                p = int(part)
                if 1 <= p <= total_pages:
                    pages.add(p - 1)
        return sorted(pages)

    def _rotate_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        in_path = files[0]
        self._ensure_pdf(in_path)
        angle = int(options.get("angle", 90)) % 360

        reader = PdfReader(in_path)
        writer = PdfWriter()

        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)

        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + f"_rotated_{angle}")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, dl_name

    def _watermark_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Add simple text watermark to each page center.
        """
        in_path = files[0]
        self._ensure_pdf(in_path)
        reader = PdfReader(in_path)
        writer = PdfWriter()

        watermark_text = options.get("watermark_text", "BlinkPDF")
        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_watermarked")

        # We will use PyMuPDF for drawing text, then rebuild PDF.
        doc = fitz.open(in_path)
        for page in doc:
            rect = page.rect
            x = rect.width / 2
            y = rect.height / 2
            page.insert_text(
                (x, y),
                watermark_text,
                color=(0.8, 0.8, 0.8),
                rotate=45,
                fontname="helv",
                fontsize=36,
                align=1,
            )
        doc.save(out_path)
        doc.close()
        return out_path, dl_name

    def _number_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Add footer page numbers to each page.
        """
        in_path = files[0]
        self._ensure_pdf(in_path)
        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_numbered")

        doc = fitz.open(in_path)
        start_at = int(options.get("start_number", 1))
        font_size = int(options.get("font_size", 10))

        for i, page in enumerate(doc):
            num = start_at + i
            rect = page.rect
            x = rect.width / 2
            y = rect.height - 20
            page.insert_text(
                (x, y),
                str(num),
                fontsize=font_size,
                fontname="helv",
                align=1,
            )

        doc.save(out_path)
        doc.close()
        return out_path, dl_name

    def _protect_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        in_path = files[0]
        self._ensure_pdf(in_path)
        password = options.get("password") or ""
        if not password:
            raise ValueError("Password is required to protect PDF.")

        reader = PdfReader(in_path)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        writer.encrypt(password)

        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_protected")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, dl_name

    def _unlock_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        in_path = files[0]
        self._ensure_pdf(in_path)
        password = options.get("password") or ""

        reader = PdfReader(in_path)
        if reader.is_encrypted:
            if not reader.decrypt(password):
                raise ValueError("Incorrect PDF password.")
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)

        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_unlocked")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, dl_name

    def _repair_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Simple repair: re-write pages with PyPDF2.
        """
        in_path = files[0]
        self._ensure_pdf(in_path)
        reader = PdfReader(in_path, strict=False)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_repaired")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, dl_name

    def _organize_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        For now, this just reorders pages if `new_order` is provided as "3,1,2".
        If not provided, it just copies the file.
        """
        in_path = files[0]
        self._ensure_pdf(in_path)
        reader = PdfReader(in_path)
        writer = PdfWriter()

        new_order = options.get("new_order", "").strip()
        if new_order:
            indices = [int(x.strip()) - 1 for x in new_order.split(",") if x.strip()]
        else:
            indices = list(range(len(reader.pages)))

        for idx in indices:
            if 0 <= idx < len(reader.pages):
                writer.add_page(reader.pages[idx])

        base_name = os.path.splitext(os.path.basename(in_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_organized")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, dl_name

    def _sign_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        If two files provided: [pdf, image_signature].
        Otherwise, just copy PDF (no server error).
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)

        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_signed")

        doc = fitz.open(pdf_path)
        sig_img_path = files[1] if len(files) > 1 else None

        for page in doc:
            if sig_img_path and os.path.exists(sig_img_path):
                rect = page.rect
                # Sign in bottom-right corner
                w = rect.width * 0.3
                h = rect.height * 0.15
                rect_sig = fitz.Rect(rect.width - w - 36, rect.height - h - 36,
                                     rect.width - 36, rect.height - 36)
                page.insert_image(rect_sig, filename=sig_img_path, keep_proportion=True)
        doc.save(out_path)
        doc.close()
        return out_path, dl_name

    def _annotate_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Minimal implementation: highlight all occurrences of a keyword if provided.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        keyword = options.get("keyword", "").strip()
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_annotated")

        if not keyword:
            # If nothing to annotate, just copy
            doc = fitz.open(pdf_path)
            doc.save(out_path)
            doc.close()
            return out_path, dl_name

        doc = fitz.open(pdf_path)
        for page in doc:
            text_instances = page.search_for(keyword)
            for inst in text_instances:
                page.add_highlight_annot(inst)
        doc.save(out_path)
        doc.close()
        return out_path, dl_name

    def _redact_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Simple redact by keyword: black out keyword occurrences.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        keyword = options.get("keyword", "").strip()
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_redacted")

        if not keyword:
            doc = fitz.open(pdf_path)
            doc.save(out_path)
            doc.close()
            return out_path, dl_name

        doc = fitz.open(pdf_path)
        for page in doc:
            areas = page.search_for(keyword)
            for rect in areas:
                page.add_redact_annot(rect, fill=(0, 0, 0))
            page.apply_redactions()
        doc.save(out_path)
        doc.close()
        return out_path, dl_name

    def _flatten_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Flatten annotations & form fields using PyMuPDF.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_flattened")

        doc = fitz.open(pdf_path)
        # Annotations & widgets are baked in by re-saving
        for page in doc:
            for annot in page.annots() or []:
                annot.set_flags(fitz.ANNOT_FLAG_PRINT)
        doc.save(out_path, deflate=True, clean=True)
        doc.close()
        return out_path, dl_name

    def _metadata_editor(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Edit basic metadata fields if provided.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_metadata")

        reader = PdfReader(pdf_path)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        meta = reader.metadata or {}
        new_meta = {}
        new_meta["/Title"] = options.get("title", meta.get("/Title", ""))
        new_meta["/Author"] = options.get("author", meta.get("/Author", ""))
        new_meta["/Subject"] = options.get("subject", meta.get("/Subject", ""))
        new_meta["/Keywords"] = options.get("keywords", meta.get("/Keywords", ""))

        writer.add_metadata(new_meta)

        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, dl_name

    def _fill_forms(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Basic form filling: options['fields'] = dict(field_name -> value).
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_filled")

        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        fields = options.get("fields") or {}

        if reader.get_fields():
            writer.append(reader)
            writer.update_page_form_field_values(writer.pages[0], fields)
        else:
            for page in reader.pages:
                writer.add_page(page)

        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, dl_name

    # ------------------------------------------------------------------
    # PAGE GEOMETRY & IMAGE-BASED OPS
    # ------------------------------------------------------------------
    def _deskew_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Basic deskew placeholder: convert pages to images and back.
        (Real Hough-based deskew needs OpenCV; here we just re-render.)
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_deskewed")

        doc = fitz.open(pdf_path)
        new_doc = fitz.open()
        zoom = 2.0
        mat = fitz.Matrix(zoom, zoom)
        for page in doc:
            pix = page.get_pixmap(matrix=mat)
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            # TODO: real deskew with extra libs; for now, keep as is.
            img_bytes = io.BytesIO()
            img.save(img_bytes, format="PNG")
            img_bytes.seek(0)
            img_pdf = fitz.open("pdf", fitz.open("png", img_bytes.read()).convert_to_pdf())
            new_doc.insert_pdf(img_pdf)

        new_doc.save(out_path)
        new_doc.close()
        doc.close()
        return out_path, dl_name

    def _crop_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Crop margins by percentage: options: top, bottom, left, right (0-40).
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_cropped")

        top_pct = float(options.get("top", 0)) / 100.0
        bottom_pct = float(options.get("bottom", 0)) / 100.0
        left_pct = float(options.get("left", 0)) / 100.0
        right_pct = float(options.get("right", 0)) / 100.0

        doc = fitz.open(pdf_path)
        for page in doc:
            rect = page.rect
            new_rect = fitz.Rect(
                rect.x0 + rect.width * left_pct,
                rect.y0 + rect.height * top_pct,
                rect.x1 - rect.width * right_pct,
                rect.y1 - rect.height * bottom_pct,
            )
            page.set_cropbox(new_rect)
        doc.save(out_path)
        doc.close()
        return out_path, dl_name

    def _resize_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Resize pages to standard sizes (A4/Letter) while scaling content.
        options["page_size"] in {"A4", "Letter"}
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        size_opt = (options.get("page_size") or "A4").upper()
        if size_opt == "LETTER":
            target_w, target_h = letter
        else:
            target_w, target_h = A4

        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + f"_{size_opt}")

        doc = fitz.open(pdf_path)
        new_doc = fitz.open()

        for page in doc:
            rect = page.rect
            scale_x = target_w / rect.width
            scale_y = target_h / rect.height
            scale = min(scale_x, scale_y)
            mat = fitz.Matrix(scale, scale)
            pix = page.get_pixmap(matrix=mat)
            img = fitz.open("png", pix.tobytes("png"))
            page_pdf = fitz.open("pdf", img.convert_to_pdf())
            page_pdf[0].set_cropbox(fitz.Rect(0, 0, target_w, target_h))
            new_doc.insert_pdf(page_pdf)

        new_doc.save(out_path)
        new_doc.close()
        doc.close()
        return out_path, dl_name

    def _background_remover(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Very basic 'background lighten' by increasing contrast on images.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_bg_clean")

        doc = fitz.open(pdf_path)
        new_doc = fitz.open()

        for page in doc:
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            # Simple threshold/whitening
            img = img.convert("L")
            img = img.point(lambda p: 255 if p > 200 else p)
            img = img.convert("RGB")

            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            img_pdf = fitz.open("pdf", fitz.open("png", buf.read()).convert_to_pdf())
            new_doc.insert_pdf(img_pdf)

        new_doc.save(out_path)
        new_doc.close()
        doc.close()
        return out_path, dl_name

    # ------------------------------------------------------------------
    # CONVERSIONS
    # ------------------------------------------------------------------
    def _pdf_to_word(self, files: List[str], options: Dict) -> Tuple[str, str]:
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name, ".docx")

        doc_pdf = fitz.open(pdf_path)
        docx = Document()
        for page in doc_pdf:
            text = page.get_text()
            docx.add_paragraph(text)
            docx.add_page_break()
        docx.save(out_path)
        doc_pdf.close()
        return out_path, dl_name
        
    def _word_to_pdf(self, input_path, output_path):
        """Convert DOCX to PDF using reportlab fallback"""

        doc = Document(input_path)
        c = canvas.Canvas(output_path)

        y = 800
        for para in doc.paragraphs:
            if not para.text.strip():
                y -= 20
                continue

            c.drawString(40, y, para.text)
            y -= 20

            if y < 40:
                c.showPage()
                y = 800

        c.save()
        return True


    def _docx_to_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Convert simple DOCX to PDF using reportlab (text-only).
        """
        docx_path = files[0]
        ext = os.path.splitext(docx_path)[1].lower()
        if ext not in [".docx"]:
            raise ValueError("Please upload a DOCX file for Word to PDF.")

        base_name = os.path.splitext(os.path.basename(docx_path))[0]
        out_path, dl_name = self._make_output_path(base_name, ".pdf")

        doc = Document(docx_path)
        c = canvas.Canvas(out_path, pagesize=A4)
        width, height = A4
        y = height - 50
        for para in doc.paragraphs:
            text = para.text
            if not text.strip():
                y -= 20
                continue
            c.drawString(50, y, text[:1200])  # simple cut
            y -= 14
            if y < 50:
                c.showPage()
                y = height - 50
        c.save()
        return out_path, dl_name

    def _pdf_to_image(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Convert each PDF page to PNG; return ZIP of images.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        zip_path, dl_name = self._make_output_path(base_name + "_images", ".zip")

        zoom = float(options.get("zoom", 2.0))
        mat = fitz.Matrix(zoom, zoom)
        doc = fitz.open(pdf_path)

        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat)
                img_name = f"{base_name}_page_{i+1}.png"
                zf.writestr(img_name, pix.tobytes("png"))

        doc.close()
        return zip_path, dl_name

    def _image_to_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Merge images into a single PDF.
        """
        base_name = "images_to_pdf"
        out_path, dl_name = self._make_output_path(base_name, ".pdf")

        pdf_bytes = io.BytesIO()
        img_list = []

        for img_path in files:
            ext = os.path.splitext(img_path)[1].lower()
            if ext not in [".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".webp"]:
                continue
            img = Image.open(img_path).convert("RGB")
            img_list.append(img)

        if not img_list:
            raise ValueError("No valid image files found for Image to PDF.")

        first, *rest = img_list
        first.save(out_path, save_all=True, append_images=rest)
        return out_path, dl_name

    def _pdf_to_excel(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Simple table-like extraction: each page as a sheet with lines of text.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name, ".xlsx")

        pdf_doc = fitz.open(pdf_path)
        wb = Workbook()
        ws0 = wb.active
        ws0.title = "Page1"

        for i, page in enumerate(pdf_doc):
            if i == 0:
                ws = ws0
            else:
                ws = wb.create_sheet(title=f"Page{i+1}")
            text = page.get_text()
            for row_idx, line in enumerate(text.splitlines(), start=1):
                ws.cell(row=row_idx, column=1, value=line)

        wb.save(out_path)
        pdf_doc.close()
        return out_path, dl_name

    def _excel_to_p
df(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Render basic Excel text into PDF using reportlab.
        """
        xlsx_path = files[0]
        ext = os.path.splitext(xlsx_path)[1].lower()
        if ext not in [".xlsx"]:
            raise ValueError("Please upload an XLSX file for Excel to PDF.")

        base_name = os.path.splitext(os.path.basename(xlsx_path))[0]
        out_path, dl_name = self._make_output_path(base_name, ".pdf")

        wb = load_workbook(xlsx_path, data_only=True)
        c = canvas.Canvas(out_path, pagesize=A4)
        width, height = A4
        y = height - 50

        for sheet in wb.worksheets:
            c.setFont("Helvetica-Bold", 12)
            c.drawString(50, y, f"Sheet: {sheet.title}")
            y -= 20
            c.setFont("Helvetica", 10)

            for row in sheet.iter_rows(values_only=True):
                line = " | ".join("" if v is None else str(v) for v in row)
                c.drawString(50, y, line[:1200])
                y -= 12
                if y < 50:
                    c.showPage()
                    y = height - 50

            c.showPage()
            y = height - 50

        c.save()
        return out_path, dl_name

    def _pdf_to_powerpoint(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Convert each page to a slide background image.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name, ".pptx")

        doc = fitz.open(pdf_path)
        prs = Presentation()
        # default 16:9 size – we just stretch images
        blank_layout = prs.slide_layouts[6]

        for page in doc:
            slide = prs.slides.add_slide(blank_layout)
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_bytes = pix.tobytes("png")
            img_stream = io.BytesIO(img_bytes)
            pic = slide.shapes.add_picture(
                img_stream,
                left=0,
                top=0,
                width=prs.slide_width,
                height=prs.slide_height,
            )

        prs.save(out_path)
        doc.close()
        return out_path, dl_name

    def _powerpoint_to_p
df(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Convert PPTX slide text into a simple PDF (not full visual render).
        """
        pptx_path = files[0]
        ext = os.path.splitext(pptx_path)[1].lower()
        if ext not in [".pptx"]:
            raise ValueError("Please upload a PPTX file for PowerPoint to PDF.")

        base_name = os.path.splitext(os.path.basename(pptx_path))[0]
        out_path, dl_name = self._make_output_path(base_name, ".pdf")

        prs = Presentation(pptx_path)
        c = canvas.Canvas(out_path, pagesize=A4)
        width, height = A4
        y = height - 50

        for i, slide in enumerate(prs.slides, start=1):
            c.setFont("Helvetica-Bold", 14)
            c.drawString(50, y, f"Slide {i}")
            y -= 24
            c.setFont("Helvetica", 11)

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    for line in shape.text.splitlines():
                        c.drawString(60, y, line[:1200])
                        y -= 14
                        if y < 50:
                            c.showPage()
                            y = height - 50
            c.showPage()
            y = height - 50

        c.save()
        return out_path, dl_name

    # ------------------------------------------------------------------
    # TEXT EXTRACTION / OCR / MEDIA
    # ------------------------------------------------------------------
    def _ocr_pdf(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        If Tesseract available: OCR each page as image.
        Otherwise, fallback to get_text().
        Returns a searchable PDF (text layer added) or a .txt if needed.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]

        if not _HAS_TESSERACT:
            # Fallback: just extract text to TXT
            out_path, dl_name = self._make_output_path(base_name + "_ocr_fallback", ".txt")
            doc = fitz.open(pdf_path)
            with open(out_path, "w", encoding="utf-8") as f:
                for page in doc:
                    f.write(page.get_text())
                    f.write("\n\n")
            doc.close()
            return out_path, dl_name

        # Real-ish OCR to TEXT file (for speed); you can later adjust to searchable PDF
        out_path, dl_name = self._make_output_path(base_name + "_ocr", ".txt")
        doc = fitz.open(pdf_path)
        with open(out_path, "w", encoding="utf-8") as f:
            for page in doc:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text = pytesseract.image_to_string(img)
                f.write(text)
                f.write("\n\n")
        doc.close()
        return out_path, dl_name

    def _extract_text(self, files: List[str], options: Dict) -> Tuple[str, str]:
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        out_path, dl_name = self._make_output_path(base_name + "_text", ".txt")

        doc = fitz.open(pdf_path)
        with open(out_path, "w", encoding="utf-8") as f:
            for page in doc:
                f.write(page.get_text())
                f.write("\n\n")
        doc.close()
        return out_path, dl_name

    def _extract_images(self, files: List[str], options: Dict) -> Tuple[str, str]:
        """
        Extract embedded images from PDF into a ZIP.
        """
        pdf_path = files[0]
        self._ensure_pdf(pdf_path)
        base_name = os.path.splitext(os.path.basename(pdf_path))[0]
        zip_path, dl_name = self._make_output_path(base_name + "_images", ".zip")

        doc = fitz.open(pdf_path)
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
            for page_index in range(len(doc)):
                page = doc[page_index]
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list, start=1):
                    xref = img_info[0]
                    base_img = doc.extract_image(xref)
                    img_bytes = base_img["image"]
                    ext = base_img["ext"]
                    img_name = f"{base_name}_p{page_index+1}_img{img_index}.{ext}"
                    zf.writestr(img_name, img_bytes)

        doc.close()
        return zip_path, dl_name

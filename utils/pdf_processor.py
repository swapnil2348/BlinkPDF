# utils/pdf_processor.py

import os
import io
import json
import tempfile
import zipfile
from typing import List, Dict, Tuple, Optional

import fitz  # PyMuPDF
from PyPDF2 import PdfReader, PdfWriter
from PIL import Image
import pytesseract
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import openpyxl
from reportlab.lib.pagesizes import A4, letter
from reportlab.pdfgen import canvas


class PDFProcessorError(Exception):
    """Custom error for PDF processing problems."""
    pass


class PDFProcessor:
    """
    Central processor for all BlinkPDF tools.

    All tools are handled via:
        process(slug: str, input_files: List[str], options: Dict[str, str])

    and return:
        (output_path: str, download_name: str, mimetype: str)
    """

    # ------------------------------------------------------------------ #
    # Public entrypoint                                                  #
    # ------------------------------------------------------------------ #

    def process(
        self,
        slug: str,
        input_files: List[str],
        options: Dict[str, str]
    ) -> Tuple[str, str, str]:
        """
        Route a tool slug to its implementation.

        slug:        tool slug like "compress-pdf", "rotate-pdf", etc.
        input_files: list of absolute file paths on disk
        options:     dict of options from the form / PRO++ front-end
        """
        slug = slug.strip().lower()
        if not input_files:
            raise PDFProcessorError("No input files provided")

        first = input_files[0]

        # -------------- Core PDF structure / utility tools -------------- #

        if slug == "merge-pdf":
            return self.merge_pdfs(input_files)

        if slug == "split-pdf":
            pages = (
                options.get("pages")
                or options.get("page_range")
                or ""
            )
            return self.split_pdf(first, pages)

        if slug == "compress-pdf":
            # Front-end sends slider 1/2/3 as string
            raw_level = options.get("compression_level", "2")
            level = self._normalize_compression_level(raw_level)
            return self.compress_pdf(first, level)

        if slug == "optimize-pdf":
            return self.optimize_pdf(first)

        if slug == "rotate-pdf":
            # PRO++ front-end uses rotation_angle; older code used angle
            angle_str = (
                options.get("rotation_angle")
                or options.get("angle")
                or "0"
            )
            try:
                angle = int(angle_str)
            except ValueError:
                angle = 0
            return self.rotate_pdf(first, angle)

        if slug == "watermark-pdf":
            text = options.get("watermark_text") or "CONFIDENTIAL"
            opacity = float(options.get("watermark_opacity", "0.15") or "0.15")
            position = options.get("watermark_position", "center")
            return self.watermark_pdf(first, text, opacity, position)

        if slug == "number-pdf":
            position = options.get("position", "bottom-right")
            return self.number_pdf(first, position)

        if slug == "protect-pdf":
            password = options.get("password") or options.get("new_password")
            if not password:
                raise PDFProcessorError("Password is required to protect PDF.")
            return self.protect_pdf(first, password)

        if slug == "unlock-pdf":
            password = options.get("password")
            if not password:
                raise PDFProcessorError("Password is required to unlock PDF.")
            return self.unlock_pdf(first, password)

        if slug == "repair-pdf":
            return self.repair_pdf(first)

        if slug == "organize-pdf":
            # PRO++:
            #   - page_order: "1,3,2" from thumbnails
            #   - deleted_pages: "4,5" for excluded pages
            # Fallback:
            #   - pages/page_range and delete_pages
            page_order_spec = (
                options.get("page_order")
                or options.get("order")
                or ""
            )
            deleted_spec = (
                options.get("deleted_pages")
                or options.get("delete_pages")
                or ""
            )
            # If front-end only used "pages" as range, still respect it
            if not page_order_spec and options.get("pages"):
                page_order_spec = options.get("pages") or ""

            return self.organize_pdf(first, page_order_spec, deleted_spec)

        if slug == "sign-pdf":
            # Expect an image signature path in options (uploaded separately via app.py)
            signature_path = options.get("signature_path")
            if not signature_path or not os.path.exists(signature_path):
                raise PDFProcessorError("Signature image not provided.")
            return self.sign_pdf(first, signature_path)

        if slug == "annotate-pdf":
            highlight_text = options.get("highlight_text")
            if not highlight_text:
                raise PDFProcessorError("highlight_text is required for annotate tool.")
            return self.annotate_pdf(first, highlight_text)

        if slug == "redact-pdf":
            redact_text = options.get("redact_text")
            if not redact_text:
                raise PDFProcessorError("redact_text is required for redact tool.")
            return self.redact_pdf(first, redact_text)

        # -------------- Conversions: PDF <-> Office / Images ------------- #

        if slug == "pdf-to-word":
            return self.pdf_to_word(first)

        if slug == "word-to-pdf":
            return self.word_to_pdf(first)

        if slug == "pdf-to-image":
            fmt = (options.get("image_format") or "png").lower()
            dpi_str = options.get("output_dpi", "150") or "150"
            try:
                dpi = int(dpi_str)
            except ValueError:
                dpi = 150
            return self.pdf_to_images(first, fmt, dpi)

        if slug == "image-to-pdf":
            return self.images_to_pdf(input_files)

        if slug == "pdf-to-excel":
            return self.pdf_to_excel(first)

        if slug == "excel-to-pdf":
            return self.excel_to_pdf(first)

        if slug == "pdf-to-powerpoint":
            return self.pdf_to_powerpoint(first)

        if slug == "powerpoint-to-pdf":
            return self.powerpoint_to_pdf(first)

        # -------------- OCR / text / images ------------------------------ #

        if slug == "ocr-pdf":
            lang = options.get("ocr_lang", "eng") or "eng"
            return self.ocr_pdf(first, lang)

        if slug == "extract-text":
            return self.extract_text(first)

        if slug == "extract-images":
            return self.extract_images(first)

        # -------------- Page geometry ------------------------------------ #

        if slug == "deskew-pdf":
            return self.deskew_pdf(first)

        if slug == "crop-pdf":
            # PRO++ crop via regions JSON has priority
            regions_json = options.get("crop_regions")
            if regions_json:
                return self.crop_pdf_regions(first, regions_json)

            # Fallback margin-based crop (points)
            top = float(options.get("crop_top", "0") or "0")
            right = float(options.get("crop_right", "0") or "0")
            bottom = float(options.get("crop_bottom", "0") or "0")
            left = float(options.get("crop_left", "0") or "0")
            return self.crop_pdf(first, top, right, bottom, left)

        if slug == "resize-pdf":
            size = options.get("page_size", "A4").upper()
            return self.resize_pdf(first, size)

        if slug == "flatten-pdf":
            return self.flatten_pdf(first)

        if slug == "metadata-editor":
            # options may contain: title, author, subject, keywords, etc.
            return self.edit_metadata(first, options)

        if slug == "fill-forms":
            data_raw = options.get("form_data_json", "{}")
            try:
                form_data = json.loads(data_raw)
            except json.JSONDecodeError:
                raise PDFProcessorError("Invalid JSON in form_data_json")
            return self.fill_forms(first, form_data)

        # -------------- Non-PDF tool: background remover ----------------- #

        if slug == "background-remover":
            # Treat first input as image, return PNG with background removed.
            return self.remove_background(first)

        raise PDFProcessorError(f"Unknown tool slug: {slug}")

    # ------------------------------------------------------------------ #
    # Helper functions                                                  #
    # ------------------------------------------------------------------ #

    def _tmp_file(self, suffix: str) -> str:
        fd, path = tempfile.mkstemp(suffix=suffix)
        os.close(fd)
        return path

    def _page_ranges_to_list(self, spec: str, total_pages: int) -> List[int]:
        """
        Convert a string like '1-3,5,7-' to zero-based page indices.

        PRO++ is also fine to send '1,3,2' (we treat it as discrete pages).
        """
        if not spec:
            return list(range(total_pages))

        pages = set()
        parts = [p.strip() for p in spec.split(",") if p.strip()]
        for part in parts:
            if "-" in part:
                start_str, end_str = part.split("-", 1)
                start = int(start_str) if start_str else 1
                end = int(end_str) if end_str else total_pages
                for p in range(start, end + 1):
                    if 1 <= p <= total_pages:
                        pages.add(p - 1)
            else:
                try:
                    p = int(part)
                except ValueError:
                    continue
                if 1 <= p <= total_pages:
                    pages.add(p - 1)
        return sorted(pages)

    def _normalize_compression_level(self, raw: str) -> str:
        """
        Map slider values 1/2/3 or words to internal levels: 'high', 'medium', 'low'
        """
        val = (raw or "").strip().lower()
        if val in {"1", "high", "hq"}:
            return "high"
        if val in {"3", "low", "smallest"}:
            return "low"
        return "medium"

    # ------------------------------------------------------------------ #
    # Implementations                                                   #
    # ------------------------------------------------------------------ #

    # ---- Merge / Split / Basic --------------------------------------- #

    def merge_pdfs(self, paths: List[str]) -> Tuple[str, str, str]:
        writer = PdfWriter()
        for path in paths:
            reader = PdfReader(path)
            for page in reader.pages:
                writer.add_page(page)
        out_path = self._tmp_file(".pdf")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, "merged.pdf", "application/pdf"

    def split_pdf(self, path: str, pages_spec: str) -> Tuple[str, str, str]:
        reader = PdfReader(path)
        total = len(reader.pages)
        indices = self._page_ranges_to_list(pages_spec, total)
        if not indices:
            indices = list(range(total))
        writer = PdfWriter()
        for i in indices:
            writer.add_page(reader.pages[i])
        out_path = self._tmp_file(".pdf")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, "split.pdf", "application/pdf"

    def compress_pdf(self, path: str, level: str) -> Tuple[str, str, str]:
        """
        Compression using PyMuPDF.

        level:
            'high'   -> prioritize quality
            'medium' -> balanced (default)
            'low'    -> prioritize smallest size
        """
        doc = fitz.open(path)
        out_path = self._tmp_file(".pdf")

        # Default parameters
        save_kwargs = {
            "deflate": True,
            "clean": True,
            "garbage": 3,
            "linear": True,
            "compress": True,
            "compress_images": True,
        }

        if level == "high":
            # keep better quality, less aggressive
            save_kwargs["garbage"] = 2
            save_kwargs["compress_images"] = True
            save_kwargs["deflate"] = True
            save_kwargs["linear"] = True
        elif level == "low":
            # push more aggressively
            save_kwargs["garbage"] = 4
            save_kwargs["deflate"] = True
            save_kwargs["linear"] = True
        else:
            # medium already set in defaults
            pass

        # image_quality parameter is supported in recent PyMuPDF
        if level == "high":
            save_kwargs["image_quality"] = 95
        elif level == "low":
            save_kwargs["image_quality"] = 60
        else:
            save_kwargs["image_quality"] = 80

        doc.save(out_path, **save_kwargs)
        doc.close()
        return out_path, "compressed.pdf", "application/pdf"

    def optimize_pdf(self, path: str) -> Tuple[str, str, str]:
        """
        Aggressive optimize similar to compress, tuned separately.
        """
        doc = fitz.open(path)
        out_path = self._tmp_file(".pdf")
        doc.save(
            out_path,
            garbage=4,
            deflate=True,
            clean=True,
            linear=True,
            compress=True,
            compress_images=True,
            image_quality=70,
            pretty=True,
        )
        doc.close()
        return out_path, "optimized.pdf", "application/pdf"

    def rotate_pdf(self, path: str, angle: int) -> Tuple[str, str, str]:
        # Normalize angle to 0,90,180,270
        angle = angle % 360
        if angle not in (0, 90, 180, 270):
            # snap to nearest
            candidates = [0, 90, 180, 270]
            angle = min(candidates, key=lambda a: abs(a - angle))

        reader = PdfReader(path)
        writer = PdfWriter()
        for page in reader.pages:
            if angle != 0:
                page.rotate(angle)
            writer.add_page(page)

        out_path = self._tmp_file(".pdf")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, "rotated.pdf", "application/pdf"

    def protect_pdf(self, path: str, password: str) -> Tuple[str, str, str]:
        reader = PdfReader(path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)
        out_path = self._tmp_file(".pdf")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, "protected.pdf", "application/pdf"

    def unlock_pdf(self, path: str, password: str) -> Tuple[str, str, str]:
        reader = PdfReader(path)
        if reader.is_encrypted:
            res = reader.decrypt(password)
            if res == 0:
                raise PDFProcessorError("Incorrect password for unlocking PDF.")
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        out_path = self._tmp_file(".pdf")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, "unlocked.pdf", "application/pdf"

    def repair_pdf(self, path: str) -> Tuple[str, str, str]:
        """
        Heuristic: reopen with PyMuPDF and resave; can fix some structural issues.
        """
        doc = fitz.open(path)
        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "repaired.pdf", "application/pdf"

    def organize_pdf(
        self,
        path: str,
        order_spec: str,
        delete_spec: str,
    ) -> Tuple[str, str, str]:
        """
        order_spec  : string like "1,3,2" or "1-3,5"
        delete_spec : pages to remove like "4,7" or "2-4"
        """
        reader = PdfReader(path)
        total = len(reader.pages)

        delete = (
            set(self._page_ranges_to_list(delete_spec, total))
            if delete_spec
            else set()
        )

        if order_spec:
            order = self._page_ranges_to_list(order_spec, total)
        else:
            order = list(range(total))

        writer = PdfWriter()
        for idx in order:
            if idx in delete:
                continue
            if 0 <= idx < total:
                writer.add_page(reader.pages[idx])

        out_path = self._tmp_file(".pdf")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, "organized.pdf", "application/pdf"

    # ---- Visual overlays: watermark, page numbers, sign, annotate, redact ---- #

    def watermark_pdf(
        self,
        path: str,
        text: str,
        opacity: float,
        position: str
    ) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        for page in doc:
            rect = page.rect
            if position == "top-left":
                point = rect.tl + (50, 50)
            elif position == "top-right":
                point = rect.tr + (-200, 50)
            elif position == "bottom-left":
                point = rect.bl + (50, -50)
            elif position == "bottom-right":
                point = rect.br + (-200, -50)
            else:  # center
                point = rect.center

            page.insert_text(
                point,
                text,
                fontsize=36,
                rotate=45,
                color=(0.7, 0.7, 0.7),
                fill_opacity=opacity,
            )
        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "watermarked.pdf", "application/pdf"

    def number_pdf(self, path: str, position: str) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        total = len(doc)
        for i, page in enumerate(doc, start=1):
            rect = page.rect
            label = f"{i}/{total}"
            if position == "top-left":
                point = rect.tl + (20, 20)
            elif position == "top-right":
                point = rect.tr + (-60, 20)
            elif position == "bottom-left":
                point = rect.bl + (20, -20)
            else:  # bottom-right
                point = rect.br + (-60, -20)
            page.insert_text(
                point,
                label,
                fontsize=10,
                color=(0, 0, 0),
            )
        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "numbered.pdf", "application/pdf"

    def sign_pdf(self, path: str, signature_image_path: str) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        sig_img = Image.open(signature_image_path).convert("RGBA")
        img_bytes_io = io.BytesIO()
        sig_img.save(img_bytes_io, format="PNG")
        img_bytes = img_bytes_io.getvalue()

        for page in doc:
            rect = page.rect
            box = fitz.Rect(
                rect.br.x - 150,
                rect.br.y - 80,
                rect.br.x - 10,
                rect.br.y - 10,
            )
            page.insert_image(box, stream=img_bytes, keep_proportion=True)

        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "signed.pdf", "application/pdf"

    def annotate_pdf(self, path: str, text: str) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        for page in doc:
            areas = page.search_for(text)
            for rect in areas:
                page.add_highlight_annot(rect)
        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "annotated.pdf", "application/pdf"

    def redact_pdf(self, path: str, text: str) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        for page in doc:
            areas = page.search_for(text)
            for rect in areas:
                page.add_redact_annot(rect, fill=(0, 0, 0))
        doc.apply_redactions()
        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "redacted.pdf", "application/pdf"

    # ---- Conversions: PDF <-> Word ------------------------------------ #

    def pdf_to_word(self, path: str) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        word_doc = Document()
        for page in doc:
            text = page.get_text()
            for line in text.splitlines():
                word_doc.add_paragraph(line)
            word_doc.add_page_break()
        doc.close()
        out_path = self._tmp_file(".docx")
        word_doc.save(out_path)
        return out_path, "converted.docx", (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    def word_to_pdf(self, path: str) -> Tuple[str, str, str]:
        """
        Simple DOCX -> PDF conversion: we take the plain text and draw it
        into a PDF using ReportLab. Layout from Word is not preserved,
        but it's a real conversion.
        """
        doc = Document(path)
        out_path = self._tmp_file(".pdf")

        c = canvas.Canvas(out_path, pagesize=A4)
        width, height = A4
        x_margin = 50
        y = height - 50

        for para in doc.paragraphs:
            text = para.text
            if not text:
                y -= 16
                if y < 50:
                    c.showPage()
                    y = height - 50
                continue
            lines = text.splitlines() or [""]
            for line in lines:
                c.drawString(x_margin, y, line[:2000])
                y -= 14
                if y < 50:
                    c.showPage()
                    y = height - 50
        c.save()

        return out_path, "converted.pdf", "application/pdf"

    # ---- Conversions: PDF <-> Images ---------------------------------- #

    def pdf_to_images(
        self,
        path: str,
        image_format: str = "png",
        dpi: int = 150
    ) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        tmpdir = tempfile.mkdtemp()
        fmt = image_format.lower()
        if fmt not in ("png", "jpg", "jpeg"):
            fmt = "png"

        image_paths = []
        zoom = dpi / 72.0
        mat = fitz.Matrix(zoom, zoom)

        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            img_path = os.path.join(tmpdir, f"page_{i + 1}.{fmt}")
            pix.save(img_path)
            image_paths.append(img_path)

        doc.close()

        out_path = self._tmp_file(".zip")
        with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as z:
            for img_path in image_paths:
                z.write(img_path, os.path.basename(img_path))

        return out_path, "images.zip", "application/zip"

    def images_to_pdf(self, paths: List[str]) -> Tuple[str, str, str]:
        """
        Combine image files into a single PDF (one image per page).
        """
        images = []
        for path in paths:
            img = Image.open(path)
            if img.mode != "RGB":
                img = img.convert("RGB")
            images.append(img)

        if not images:
            raise PDFProcessorError("No valid images provided.")

        out_path = self._tmp_file(".pdf")
        first, *rest = images
        first.save(out_path, save_all=True, append_images=rest)
        return out_path, "images.pdf", "application/pdf"

    # ---- Conversions: PDF <-> Excel ----------------------------------- #

    def pdf_to_excel(self, path: str) -> Tuple[str, str, str]:
        """
        Very simple PDF->Excel: extract text page by page, write each line
        into successive rows. Real export but not a full table detector.
        """
        doc = fitz.open(path)
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "PDF Text"

        row = 1
        for page_index, page in enumerate(doc, start=1):
            text = page.get_text()
            ws.cell(row=row, column=1, value=f"=== Page {page_index} ===")
            row += 1
            for line in text.splitlines():
                ws.cell(row=row, column=1, value=line)
                row += 1
            row += 1  # blank line between pages

        doc.close()
        out_path = self._tmp_file(".xlsx")
        wb.save(out_path)
        return out_path, "converted.xlsx", (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        )

    def excel_to_pdf(self, path: str) -> Tuple[str, str, str]:
        """
        Simple Excel->PDF: we draw each row into the PDF using ReportLab.
        Layout/styling is not preserved, but all cell values are exported.
        """
        wb = openpyxl.load_workbook(path, data_only=True)
        sheet = wb.active

        out_path = self._tmp_file(".pdf")
        c = canvas.Canvas(out_path, pagesize=letter)
        width, height = letter
        x_margin = 40
        y = height - 40

        for row_cells in sheet.iter_rows():
            values = [
                str(cell.value) if cell.value is not None else ""
                for cell in row_cells
            ]
            line = " | ".join(values)
            if not line.strip():
                y -= 14
            else:
                c.drawString(x_margin, y, line[:2000])
                y -= 14

            if y < 40:
                c.showPage()
                y = height - 40

        c.save()
        return out_path, "converted.pdf", "application/pdf"

    # ---- Conversions: PDF <-> PowerPoint ------------------------------- #

    def pdf_to_powerpoint(self, path: str) -> Tuple[str, str, str]:
        """
        PDF->PPTX: each page rendered as an image and inserted as a full-slide image.
        """
        doc = fitz.open(path)
        prs = Presentation()

        blank_layout = prs.slide_layouts[6]  # blank
        zoom = 150 / 72.0
        mat = fitz.Matrix(zoom, zoom)

        for page in doc:
            pix = page.get_pixmap(matrix=mat)
            img_path = self._tmp_file(".png")
            pix.save(img_path)

            slide = prs.slides.add_slide(blank_layout)
            slide_width = prs.slide_width
            slide_height = prs.slide_height

            slide.shapes.add_picture(
                img_path,
                0,
                0,
                width=slide_width,
                height=slide_height,
            )

        doc.close()
        out_path = self._tmp_file(".pptx")
        prs.save(out_path)
        return out_path, "converted.pptx", (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    def powerpoint_to_pdf(self, path: str) -> Tuple[str, str, str]:
        """
        PPTX->PDF without a rendering engine: export text content of each slide
        into a PDF using ReportLab. Not layout-accurate but real conversion.
        """
        prs = Presentation(path)
        out_path = self._tmp_file(".pdf")

        c = canvas.Canvas(out_path, pagesize=letter)
        width, height = letter
        x_margin = 40
        y = height - 40

        for slide_index, slide in enumerate(prs.slides, start=1):
            c.setFont("Helvetica-Bold", 12)
            c.drawString(x_margin, y, f"Slide {slide_index}")
            y -= 16
            c.setFont("Helvetica", 10)

            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text or ""
                for line in text.splitlines():
                    c.drawString(x_margin, y, line[:2000])
                    y -= 12
                    if y < 40:
                        c.showPage()
                        y = height - 40
            y -= 20
            if y < 40:
                c.showPage()
                y = height - 40

        c.save()
        return out_path, "converted.pdf", "application/pdf"

    # ---- OCR / text / images ------------------------------------------ #

    def ocr_pdf(self, path: str, lang: str = "eng") -> Tuple[str, str, str]:
        """
        Perform OCR on each page and create a new searchable PDF.
        Requires Tesseract installed on system. If not installed,
        pytesseract will raise an error which we convert to a clean message.
        """
        try:
            doc = fitz.open(path)
            out_path = self._tmp_file(".pdf")

            pdf_bytes = b""
            for page in doc:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                pdf_bytes += pytesseract.image_to_pdf_or_hocr(
                    img,
                    lang=lang,
                    extension="pdf",
                )

            doc.close()
            with open(out_path, "wb") as f:
                f.write(pdf_bytes)
            return out_path, "ocr.pdf", "application/pdf"
        except pytesseract.TesseractNotFoundError:
            raise PDFProcessorError(
                "Tesseract OCR binary not found on server. "
                "Install it to enable OCR functionality."
            )

    def extract_text(self, path: str) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        texts = []
        for i, page in enumerate(doc, start=1):
            texts.append(f"=== Page {i} ===\n")
            texts.append(page.get_text())
            texts.append("\n\n")
        doc.close()
        out_path = self._tmp_file(".txt")
        with open(out_path, "w", encoding="utf-8") as f:
            f.write("".join(texts))
        return out_path, "extracted.txt", "text/plain"

    def extract_images(self, path: str) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        tmpdir = tempfile.mkdtemp()
        img_paths = []

        for page_index in range(len(doc)):
            for img_index, img in enumerate(
                doc.get_page_images(page_index), start=1
            ):
                xref = img[0]
                base = doc.extract_image(xref)
                img_bytes = base["image"]
                img_ext = base["ext"]
                img_path = os.path.join(
                    tmpdir,
                    f"page{page_index + 1}_img{img_index}.{img_ext}",
                )
                with open(img_path, "wb") as f:
                    f.write(img_bytes)
                img_paths.append(img_path)

        doc.close()

        if not img_paths:
            raise PDFProcessorError("No embedded images found in PDF.")

        out_path = self._tmp_file(".zip")
        with zipfile.ZipFile(out_path, "w", zipfile.ZIP_DEFLATED) as z:
            for p in img_paths:
                z.write(p, os.path.basename(p))
        return out_path, "images.zip", "application/zip"

    # ---- Page geometry: deskew / crop / resize / flatten -------------- #

    def deskew_pdf(self, path: str) -> Tuple[str, str, str]:
        """
        Use Tesseract's OSD to detect rotation for each page image
        and rotate it back. Real but depends on Tesseract.
        """
        try:
            doc = fitz.open(path)
            new_doc = fitz.open()

            for page in doc:
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                osd = pytesseract.image_to_osd(img)
                angle = 0
                for line in osd.splitlines():
                    if "Rotate:" in line:
                        try:
                            angle = int(line.split(":")[1].strip())
                        except ValueError:
                            angle = 0
                        break
                if angle != 0:
                    img = img.rotate(-angle, expand=True)

                buf = io.BytesIO()
                img.save(buf, format="PNG")
                buf.seek(0)
                new_page_rect = fitz.Rect(0, 0, img.width, img.height)
                new_page = new_doc.new_page(
                    width=img.width,
                    height=img.height,
                )
                new_page.insert_image(new_page_rect, stream=buf.getvalue())

            doc.close()
            out_path = self._tmp_file(".pdf")
            new_doc.save(out_path)
            new_doc.close()
            return out_path, "deskewed.pdf", "application/pdf"
        except pytesseract.TesseractNotFoundError:
            raise PDFProcessorError(
                "Tesseract OCR binary not found on server. "
                "Install it to enable deskew functionality."
            )

    def crop_pdf(
        self,
        path: str,
        top: float,
        right: float,
        bottom: float,
        left: float
    ) -> Tuple[str, str, str]:
        """
        Simple margin crop in PDF points.
        """
        doc = fitz.open(path)
        for page in doc:
            r = page.rect
            new_rect = fitz.Rect(
                r.x0 + left,
                r.y0 + top,
                r.x1 - right,
                r.y1 - bottom,
            )
            if new_rect.width <= 0 or new_rect.height <= 0:
                continue
            page.set_cropbox(new_rect)
        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "cropped.pdf", "application/pdf"

    def crop_pdf_regions(
        self,
        path: str,
        regions_json: str
    ) -> Tuple[str, str, str]:
        """
        PRO++ crop: regions_json is a JSON object like:
        {
          "1": { "x": 0.1, "y": 0.1, "width": 0.8, "height": 0.8 },
          "2": { ... },
          ...
        }
        where values are relative [0,1] to the page.
        """
        try:
            data = json.loads(regions_json) if regions_json else {}
        except json.JSONDecodeError:
            raise PDFProcessorError("Invalid crop_regions JSON.")

        if not isinstance(data, dict):
            raise PDFProcessorError("crop_regions JSON must be an object.")

        doc = fitz.open(path)
        num_pages = len(doc)

        for page_index in range(num_pages):
            page_no = page_index + 1
            region = data.get(str(page_no)) or data.get(page_no)
            if not region:
                # if no region, leave page as-is
                continue

            try:
                x = float(region.get("x", 0.0))
                y = float(region.get("y", 0.0))
                width_rel = float(region.get("width", 1.0))
                height_rel = float(region.get("height", 1.0))
            except (TypeError, ValueError):
                continue

            r = doc[page_index].rect
            x = max(0.0, min(1.0, x))
            y = max(0.0, min(1.0, y))
            width_rel = max(0.0, min(1.0, width_rel))
            height_rel = max(0.0, min(1.0, height_rel))

            x1 = r.x0 + r.width * x
            y1 = r.y0 + r.height * y
            x2 = r.x0 + r.width * (x + width_rel)
            y2 = r.y0 + r.height * (y + height_rel)

            new_rect = fitz.Rect(x1, y1, x2, y2)
            if new_rect.width <= 0 or new_rect.height <= 0:
                continue

            page = doc[page_index]
            page.set_cropbox(new_rect)

        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "cropped.pdf", "application/pdf"

    def resize_pdf(self, path: str, size: str) -> Tuple[str, str, str]:
        size = size.upper()
        if size == "LETTER":
            page_size = letter
        else:
            page_size = A4

        target_w, target_h = page_size
        doc = fitz.open(path)
        new_doc = fitz.open()

        for page in doc:
            pix = page.get_pixmap()
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)

            new_page = new_doc.new_page(width=target_w, height=target_h)
            rect = fitz.Rect(0, 0, target_w, target_h)
            new_page.insert_image(rect, stream=buf.getvalue(), keep_proportion=True)

        doc.close()
        out_path = self._tmp_file(".pdf")
        new_doc.save(out_path)
        new_doc.close()
        return out_path, "resized.pdf", "application/pdf"

    def flatten_pdf(self, path: str) -> Tuple[str, str, str]:
        doc = fitz.open(path)
        for page in doc:
            page.flatten_annotations()
        out_path = self._tmp_file(".pdf")
        doc.save(out_path)
        doc.close()
        return out_path, "flattened.pdf", "application/pdf"

    # ---- Metadata / Forms ---------------------------------------------- #

    def edit_metadata(self, path: str, meta: Dict[str, str]) -> Tuple[str, str, str]:
        reader = PdfReader(path)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        cleaned = {
            k: v
            for k, v in meta.items()
            if v and k in {"title", "author", "subject", "keywords", "creator", "producer"}
        }
        if cleaned:
            writer.add_metadata({
                f"/{k.capitalize()}": v for k, v in cleaned.items()
            })

        out_path = self._tmp_file(".pdf")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, "metadata.pdf", "application/pdf"

    def fill_forms(self, path: str, form_data: Dict[str, str]) -> Tuple[str, str, str]:
        reader = PdfReader(path)
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        if writer.pages:
            writer.update_page_form_field_values(writer.pages[0], form_data or {})

        out_path = self._tmp_file(".pdf")
        with open(out_path, "wb") as f:
            writer.write(f)
        return out_path, "filled.pdf", "application/pdf"

    # ---- Background Remover (image tool) -------------------------------- #

    def remove_background(self, path: str) -> Tuple[str, str, str]:
        """
        Image background remover.

        If 'rembg' library is installed, we use it for high-quality removal.
        If not, we raise a clear error instead of faking the result.
        """
        ext = os.path.splitext(path)[1].lower()
        if ext not in [".png", ".jpg", ".jpeg", ".webp"]:
            raise PDFProcessorError("Background remover expects an image file.")

        try:
            from rembg import remove  # type: ignore
        except ImportError:
            raise PDFProcessorError(
                "Background remover requires the 'rembg' package. "
                "Add 'rembg' to requirements.txt and deploy again."
            )

        with open(path, "rb") as f:
            input_bytes = f.read()
        output_bytes = remove(input_bytes)

        out_path = self._tmp_file(".png")
        with open(out_path, "wb") as f:
            f.write(output_bytes)

        return out_path, "no-background.png", "image/png"
